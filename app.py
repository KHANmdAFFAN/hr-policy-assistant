import streamlit as st
from pathlib import Path
from pptx import Presentation

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain_community.vectorstores import FAISS

# ===============================
# PAGE CONFIG
# ===============================

st.set_page_config(
    page_title="HR Policy Assistant",
    page_icon="🤖",
    layout="wide"
)

st.title("🤖 HR Policy Assistant")
st.write("Ask questions about HR policies.")

# ===============================
# GROQ API KEY
# ===============================

groq_api = st.secrets["GROQ_API_KEY"]

# ===============================
# LOAD DOCUMENTS
# ===============================

@st.cache_resource
def load_documents():

    data_dir = Path(__file__).resolve().parent / "data"

    documents = []

    pptx_files = list(data_dir.glob("*.pptx"))

    if not pptx_files:
        raise ValueError("No PPTX files found in data folder.")

    for pptx_file in pptx_files:

        prs = Presentation(str(pptx_file))

        slide_text = []

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:
                        slide_text.append(text)

        documents.append(
            Document(
                page_content="\n".join(slide_text),
                metadata={"source": pptx_file.name}
            )
        )

    return documents


# ===============================
# SPLIT DOCUMENTS
# ===============================

@st.cache_resource
def split_documents(documents):

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )

    return splitter.split_documents(documents)


# ===============================
# VECTOR STORE
# ===============================

@st.cache_resource
def create_vector_store(docs):

    embeddings = HuggingFaceEmbeddings(
        model_name="BAAI/bge-small-en-v1.5"
    )

    vector_store = FAISS.from_documents(
        docs,
        embeddings
    )

    return vector_store


# ===============================
# LOAD LLM
# ===============================

@st.cache_resource
def load_llm():

    return ChatGroq(
        groq_api_key=groq_api,
        model_name="llama-3.3-70b-versatile"
    )


# ===============================
# RAG RESPONSE
# ===============================

def hr_agent(query):

    docs = vector_store.similarity_search(
        query,
        k=5
    )

    context = "\n\n".join(
        [doc.page_content for doc in docs]
    )

    prompt = f"""
You are an HR Policy Assistant.

Answer only from the provided HR policy context.

Context:
{context}

Question:
{query}

Answer:
"""

    response = llm.invoke(prompt)

    return response.content


# ===============================
# INITIALIZE SYSTEM
# ===============================

try:

    with st.spinner("Loading HR Policies..."):

        documents = load_documents()

        docs = split_documents(documents)

        vector_store = create_vector_store(docs)

        llm = load_llm()

    st.success("System Ready ✅")

except Exception as e:

    st.error(f"Initialization failed: {str(e)}")
    st.stop()


# ===============================
# USER INPUT
# ===============================

query = st.text_input(
    "Enter your HR question:"
)

if st.button("Ask"):

    if query.strip():

        with st.spinner("Thinking..."):

            answer = hr_agent(query)

        st.subheader("Answer")

        st.write(answer)

    else:

        st.warning(
            "Please enter a question."
        )


# ===============================
# SIDEBAR
# ===============================

st.sidebar.title("About")

st.sidebar.info(
    """
    HR Policy Assistant

    Powered by:
    - LangChain
    - FAISS
    - HuggingFace Embeddings
    - Groq Llama 3.3 70B
    """
)